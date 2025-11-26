
import PyPDF2
from pptx import Presentation

def extract_text_from_file(file):
    if file.type == "application/pdf":
        reader = PyPDF2.PdfReader(file)
        return "\n".join(page.extract_text() for page in reader.pages)
    elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        prs = Presentation(file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    else:
        return None
